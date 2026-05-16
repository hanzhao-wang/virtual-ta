from __future__ import annotations

import csv
import json
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from .common import normalize_text, slugify


TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".py", ".json", ".yaml", ".yml"}
OFFICE_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}
DATA_EXTENSIONS = {".csv", ".html", ".htm", ".ipynb", ".tex"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
ARCHIVE_EXTENSIONS = {".zip"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | DATA_EXTENSIONS | IMAGE_EXTENSIONS | ARCHIVE_EXTENSIONS
LEGACY_EXTENSIONS = {".ppt", ".doc", ".xls"}
ALL_HANDLED_EXTENSIONS = SUPPORTED_EXTENSIONS | LEGACY_EXTENSIONS


def make_segment(
    text: str,
    *,
    locator_type: str = "document",
    locator: str = "document",
    title: str = "",
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return {
        "text": normalize_text(text),
        "locator_type": locator_type,
        "locator": str(locator),
        "title": title.strip(),
        "metadata": metadata or {},
    }


def combine_segments(segments: list[dict[str, Any]]) -> str:
    parts = []
    for segment in segments:
        title = segment.get("title") or f"{segment.get('locator_type', 'document')} {segment.get('locator', '')}"
        text = segment.get("text", "")
        if text:
            parts.append(f"## {title}\n{text}")
    return normalize_text("\n\n".join(parts))


def caption_path_for_image(path: Path, caption_root: Path) -> Path:
    return caption_root / f"{slugify(str(path.resolve()))}.md"


def is_caption_stub(text: str) -> bool:
    stripped = text.strip()
    return not stripped or stripped.startswith("# Caption needed") or "Caption needed" in stripped[:120]


def extract_text(path: Path) -> tuple[str, dict[str, Any]]:
    """Backward-compatible API for older callers."""
    segments, meta = extract_material(path)
    return combine_segments(segments), meta


def extract_material(
    path: Path,
    *,
    caption_root: Path | None = None,
    config: dict[str, Any] | None = None,
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    suffix = path.suffix.lower()
    config = config or {}

    if suffix in TEXT_EXTENSIONS:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [make_segment(text, title=path.stem)], {"method": "plain_text", "status": "indexed", "warnings": []}

    if suffix == ".csv":
        return extract_csv_segments(path, max_rows=int(config.get("csv_max_rows", 30)))

    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="ignore")
        return [make_segment(extract_html(html), title=path.stem)], {"method": "html", "status": "indexed", "warnings": []}

    if suffix == ".ipynb":
        return extract_ipynb_segments(path)

    if suffix == ".tex":
        return extract_tex_segments(path)

    if suffix == ".pdf":
        return extract_pdf_segments(path, low_text_chars=int(config.get("pdf_low_text_chars_per_page", 40)))

    if suffix == ".pptx":
        return extract_pptx_segments(path)

    if suffix == ".docx":
        return extract_docx_segments(path)

    if suffix == ".xlsx":
        return extract_xlsx_segments(path, max_rows=int(config.get("xlsx_max_rows_per_sheet", 30)))

    if suffix in IMAGE_EXTENSIONS:
        return extract_image_caption(path, caption_root=caption_root, stub_template=config.get("caption_stub_template"))

    if suffix == ".zip":
        return extract_zip_segments(path, caption_root=caption_root, config=config)

    if suffix in LEGACY_EXTENSIONS:
        converted = maybe_convert_office(path)
        if converted is None:
            warning = (
                f"Legacy Office file `{path.name}` needs LibreOffice (`soffice` or `libreoffice`) "
                "or manual conversion before it can be indexed."
            )
            return [], {"method": "legacy_office", "status": "unsupported", "warnings": [warning]}
        segments, meta = extract_material(converted, caption_root=caption_root, config=config)
        meta = {**meta, "method": f"converted_{meta.get('method', 'unknown')}", "status": "converted"}
        meta.setdefault("warnings", []).append(f"Converted from legacy Office file `{path.name}`.")
        return segments, meta

    warning = f"Unsupported file type `{path.suffix}`. Convert it to PDF, PPTX, DOCX, TXT, Markdown, CSV, XLSX, TeX, or HTML."
    return [], {"method": "unsupported", "status": "unsupported", "warnings": [warning]}


def extract_csv_segments(path: Path, *, max_rows: int = 30) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    rows: list[list[str]] = []
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row_idx, row in enumerate(reader, start=1):
            if row_idx > max_rows + 1:
                break
            rows.append([cell.strip() for cell in row])

    if not rows:
        return [], {"method": "csv", "status": "failed", "warnings": [f"No readable rows in `{path.name}`."]}

    header = rows[0]
    sample_rows = rows[1:]
    lines = [
        f"CSV file: {path.name}",
        f"Columns ({len(header)}): " + " | ".join(header),
        f"Sample rows indexed: {len(sample_rows)}",
        "",
    ]
    for row_idx, row in enumerate(sample_rows, start=1):
        lines.append(f"Row {row_idx}: " + " | ".join(row))

    segment = make_segment(
        "\n".join(lines),
        locator_type="rows",
        locator=f"header and first {len(sample_rows)} rows",
        title=f"{path.stem} data summary",
        metadata={"columns": header, "sample_rows": len(sample_rows)},
    )
    return [segment], {"method": "csv_summary", "status": "indexed", "warnings": []}


def extract_csv(path: Path) -> str:
    segments, _ = extract_csv_segments(path)
    return combine_segments(segments)


def extract_html(html: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ModuleNotFoundError:
        text = re.sub(r"<[^>]+>", " ", html)
        return normalize_text(text)

    soup = BeautifulSoup(html, "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    text = soup.get_text(separator="\n")
    return normalize_text(text)


def extract_ipynb_segments(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    notebook = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    segments: list[dict[str, Any]] = []

    for idx, cell in enumerate(notebook.get("cells", []), start=1):
        cell_type = cell.get("cell_type", "cell")
        source = "".join(cell.get("source", [])).strip()
        outputs: list[str] = []
        for output in cell.get("outputs", []):
            text = output.get("text")
            if isinstance(text, list):
                outputs.append("".join(text))
            elif isinstance(text, str):
                outputs.append(text)

            data = output.get("data", {})
            plain = data.get("text/plain")
            if isinstance(plain, list):
                outputs.append("".join(plain))
            elif isinstance(plain, str):
                outputs.append(plain)

        parts = [source]
        if outputs:
            parts.append("Output:\n" + "\n".join(outputs[:3]))
        text = normalize_text("\n\n".join(part for part in parts if part.strip()))
        if text:
            segments.append(
                make_segment(
                    text,
                    locator_type="cell",
                    locator=str(idx),
                    title=f"{cell_type.title()} cell {idx}",
                    metadata={"cell_type": cell_type},
                )
            )

    status = "indexed" if segments else "failed"
    warnings = [] if segments else [f"No readable cells in `{path.name}`."]
    return segments, {"method": "ipynb", "status": status, "warnings": warnings, "cell_count": len(segments)}


def extract_tex_segments(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    cleaned = re.sub(r"%.*", "", raw)
    markers = list(re.finditer(r"\\(?:section|subsection|subsubsection)\*?\{([^}]+)\}|\\item\b", cleaned))

    if not markers:
        text = latex_to_plain_text(cleaned)
        return [make_segment(text, title=path.stem)], {"method": "tex", "status": "indexed", "warnings": []}

    segments: list[dict[str, Any]] = []
    question_counter = 0
    current_section = path.stem
    for marker_idx, marker in enumerate(markers):
        start = marker.end()
        end = markers[marker_idx + 1].start() if marker_idx + 1 < len(markers) else len(cleaned)
        body = latex_to_plain_text(cleaned[start:end])
        if not body:
            if marker.group(1):
                current_section = marker.group(1).strip()
            continue

        if marker.group(1):
            current_section = marker.group(1).strip()
            locator_type = "section"
            locator = current_section
            title = current_section
        else:
            question_counter += 1
            locator_type = "question"
            locator = str(question_counter)
            title = f"{current_section} question {question_counter}"
        segments.append(make_segment(body, locator_type=locator_type, locator=locator, title=title))

    status = "indexed" if segments else "failed"
    warnings = [] if segments else [f"No readable TeX content in `{path.name}`."]
    return segments, {"method": "tex", "status": status, "warnings": warnings, "segment_count": len(segments)}


def latex_to_plain_text(text: str) -> str:
    text = re.sub(r"\\begin\{[^}]+\}|\\end\{[^}]+\}", "\n", text)
    text = re.sub(r"\\(?:textbf|textit|emph|underline|section|subsection|subsubsection)\*?\{([^}]*)\}", r"\1", text)
    text = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^}]*\})?", " ", text)
    text = text.replace("\\$", "$").replace("\\%", "%").replace("\\_", "_")
    text = text.replace("{", " ").replace("}", " ")
    return normalize_text(text)


def extract_pdf_segments(path: Path, *, low_text_chars: int = 40) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ModuleNotFoundError as exc:
        raise RuntimeError("Missing dependency pypdf. Install requirements.txt.") from exc

    reader = PdfReader(str(path))
    segments: list[dict[str, Any]] = []
    empty_pages = 0
    low_text_pages = 0
    for idx, page in enumerate(reader.pages, start=1):
        page_text = normalize_text(page.extract_text() or "")
        if not page_text:
            empty_pages += 1
            continue
        if len(page_text) < low_text_chars:
            low_text_pages += 1
        segments.append(make_segment(page_text, locator_type="page", locator=str(idx), title=f"Page {idx}"))

    warnings: list[str] = []
    if empty_pages:
        warnings.append(f"{empty_pages} PDF page(s) had no extractable text; scanned/image-only pages may need captions or OCR.")
    if low_text_pages:
        warnings.append(f"{low_text_pages} PDF page(s) had very little text; check extraction quality.")

    if not segments:
        return [], {"method": "pdf", "status": "caption_needed", "warnings": warnings or ["No extractable PDF text found."]}
    status = "partial" if warnings else "indexed"
    return segments, {"method": "pdf", "status": status, "warnings": warnings, "page_count": len(reader.pages)}


def extract_pptx_segments(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError("Missing dependency python-pptx. Install requirements.txt.") from exc

    presentation = Presentation(str(path))
    segments: list[dict[str, Any]] = []
    empty_slides = 0

    for idx, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                text_parts.append(text.strip())
        slide_text = normalize_text("\n".join(text_parts))
        if slide_text:
            segments.append(make_segment(slide_text, locator_type="slide", locator=str(idx), title=f"Slide {idx}"))
        else:
            empty_slides += 1

    warnings = []
    if empty_slides:
        warnings.append(f"{empty_slides} slide(s) had no extractable text; image-only slides may need captions.")
    if not segments:
        return [], {"method": "pptx", "status": "caption_needed", "warnings": warnings or ["No extractable slide text found."]}
    status = "partial" if warnings else "indexed"
    return segments, {"method": "pptx", "status": status, "warnings": warnings, "slide_count": len(presentation.slides)}


def extract_docx_segments(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    try:
        from docx import Document
    except ModuleNotFoundError as exc:
        raise RuntimeError("Missing dependency python-docx. Install requirements.txt.") from exc

    document = Document(str(path))
    segments: list[dict[str, Any]] = []
    current_heading = "document"
    current_parts: list[str] = []

    def flush() -> None:
        nonlocal current_parts
        text = normalize_text("\n\n".join(current_parts))
        if text:
            segments.append(
                make_segment(text, locator_type="heading", locator=current_heading, title=current_heading)
            )
        current_parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "")
        if style_name.lower().startswith("heading"):
            flush()
            current_heading = text
        else:
            current_parts.append(text)
    flush()

    for table_idx, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                rows.append(row_text)
        if rows:
            segments.append(
                make_segment(
                    "\n".join(rows),
                    locator_type="table",
                    locator=str(table_idx),
                    title=f"Table {table_idx}",
                )
            )

    status = "indexed" if segments else "failed"
    warnings = [] if segments else [f"No readable Word content in `{path.name}`."]
    return segments, {"method": "docx", "status": status, "warnings": warnings, "segment_count": len(segments)}


def extract_xlsx_segments(path: Path, *, max_rows: int = 30) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError as exc:
        raise RuntimeError("Missing dependency openpyxl. Install requirements.txt.") from exc

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    segments: list[dict[str, Any]] = []

    for sheet in workbook.worksheets:
        rows = []
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_idx > max_rows:
                break
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append(values)

        if not rows:
            continue
        width = max(len(row) for row in rows)
        first_row = rows[0]
        lines = [
            f"Workbook: {path.name}",
            f"Sheet: {sheet.title}",
            f"Dimensions: {sheet.max_row} rows x {sheet.max_column} columns",
            "First row / likely columns: " + " | ".join(first_row),
            f"Sample rows indexed: {len(rows)}",
            "",
        ]
        for row_idx, row in enumerate(rows[:max_rows], start=1):
            padded = row + [""] * (width - len(row))
            lines.append(f"Row {row_idx}: " + " | ".join(padded))
        segments.append(
            make_segment(
                "\n".join(lines),
                locator_type="sheet",
                locator=sheet.title,
                title=f"{sheet.title} sheet summary",
                metadata={"sheet": sheet.title, "rows": sheet.max_row, "columns": sheet.max_column},
            )
        )

    workbook.close()
    status = "indexed" if segments else "failed"
    warnings = [] if segments else [f"No readable workbook content in `{path.name}`."]
    return segments, {"method": "xlsx_summary", "status": status, "warnings": warnings, "sheet_count": len(segments)}


def extract_image_caption(
    path: Path,
    *,
    caption_root: Path | None,
    stub_template: str | None = None,
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    if caption_root is None:
        warning = f"Image `{path.name}` needs a caption, but no caption root is configured."
        return [], {"method": "image_caption", "status": "caption_needed", "warnings": [warning]}

    caption_root.mkdir(parents=True, exist_ok=True)
    caption_path = caption_path_for_image(path, caption_root)
    if not caption_path.exists():
        template = stub_template or "Describe the image content, labels, axes, and course concept."
        caption_path.write_text(
            f"# Caption needed for {path}\n\n{template}\n",
            encoding="utf-8",
        )
        warning = f"Caption needed for image `{path.name}`. Fill `{caption_path.name}` under the configured caption root and rerun indexing."
        return [], {
            "method": "image_caption",
            "status": "caption_needed",
            "warnings": [warning],
            "caption_path": str(caption_path),
        }

    caption = normalize_text(caption_path.read_text(encoding="utf-8", errors="ignore"))
    if is_caption_stub(caption):
        warning = f"Caption file `{caption_path.name}` is still a stub."
        return [], {
            "method": "image_caption",
            "status": "caption_needed",
            "warnings": [warning],
            "caption_path": str(caption_path),
        }

    segment = make_segment(
        caption,
        locator_type="caption",
        locator=caption_path.name,
        title=f"Caption for {path.name}",
        metadata={"caption_path": str(caption_path)},
    )
    return [segment], {"method": "image_caption", "status": "indexed", "warnings": [], "caption_path": str(caption_path)}


def extract_zip_segments(
    path: Path,
    *,
    caption_root: Path | None,
    config: dict[str, Any],
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    max_members = int(config.get("archive_max_members", 200))
    max_member_bytes = int(config.get("archive_max_member_mb", 25)) * 1024 * 1024
    warnings: list[str] = []
    segments: list[dict[str, Any]] = []

    with zipfile.ZipFile(path) as archive:
        infos = [info for info in archive.infolist() if not info.is_dir()]
        manifest_lines = [
            f"Archive: {path.name}",
            f"File count: {len(infos)}",
            "",
            "Members:",
        ]
        for info in infos[:max_members]:
            manifest_lines.append(f"- {info.filename} ({info.file_size} bytes)")
        if len(infos) > max_members:
            warnings.append(f"Archive has {len(infos)} files; only the first {max_members} were listed/indexed.")

        segments.append(
            make_segment(
                "\n".join(manifest_lines),
                locator_type="archive",
                locator="manifest",
                title=f"{path.name} archive manifest",
            )
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_root = Path(temp_dir)
            for info in infos[:max_members]:
                suffix = Path(info.filename).suffix.lower()
                if suffix not in SUPPORTED_EXTENSIONS - IMAGE_EXTENSIONS - ARCHIVE_EXTENSIONS:
                    continue
                if info.file_size > max_member_bytes:
                    warnings.append(f"Skipped `{info.filename}` inside archive because it exceeds the configured size limit.")
                    continue
                target = temp_root / f"{slugify(info.filename)}{suffix}"
                target.write_bytes(archive.read(info))
                try:
                    member_segments, member_meta = extract_material(target, caption_root=caption_root, config=config)
                except Exception as exc:  # noqa: BLE001
                    warnings.append(f"Failed to extract `{info.filename}` inside archive: {exc}")
                    continue
                warnings.extend(member_meta.get("warnings", []))
                for segment in member_segments:
                    locator_type = segment.get("locator_type", "document")
                    locator = segment.get("locator", "document")
                    segments.append(
                        {
                            **segment,
                            "locator": f"{info.filename} :: {locator_type} {locator}",
                            "metadata": {**segment.get("metadata", {}), "archive_member": info.filename},
                        }
                    )

    status = "partial" if warnings else "indexed"
    return segments, {"method": "zip", "status": status, "warnings": warnings, "member_count": len(segments)}


def maybe_convert_office(path: Path) -> Path | None:
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if libreoffice is None:
        return None

    target_ext = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}.get(path.suffix.lower())
    if target_ext is None:
        return None

    persistent_dir = Path(tempfile.mkdtemp(prefix="virtual_ta_converted_"))
    try:
        subprocess.run(
            [
                libreoffice,
                "--headless",
                "--convert-to",
                target_ext.lstrip("."),
                str(path),
                "--outdir",
                str(persistent_dir),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
    except subprocess.CalledProcessError:
        shutil.rmtree(persistent_dir, ignore_errors=True)
        return None

    converted = persistent_dir / f"{path.stem}{target_ext}"
    if converted.exists():
        return converted
    shutil.rmtree(persistent_dir, ignore_errors=True)
    return None
