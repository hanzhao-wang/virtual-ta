from __future__ import annotations

import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tools"))

from lib.extract import caption_path_for_image, extract_material  # noqa: E402


def test_blank_pdf_is_reported_as_caption_needed(tmp_path: Path) -> None:
    from pypdf import PdfWriter

    pdf_path = tmp_path / "scanned.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with pdf_path.open("wb") as handle:
        writer.write(handle)

    segments, meta = extract_material(pdf_path)

    assert segments == []
    assert meta["status"] == "caption_needed"
    assert any("no extractable text" in warning for warning in meta["warnings"])


def test_pptx_extraction_preserves_slide_locator(tmp_path: Path) -> None:
    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Model Evaluation"
    slide.placeholders[1].text = "Validation sets support model selection."
    presentation.save(pptx_path)

    segments, meta = extract_material(pptx_path)

    assert meta["status"] == "indexed"
    assert segments[0]["locator_type"] == "slide"
    assert segments[0]["locator"] == "1"
    assert "Validation sets" in segments[0]["text"]


def test_docx_extraction_preserves_heading_locator(tmp_path: Path) -> None:
    from docx import Document

    docx_path = tmp_path / "notes.docx"
    document = Document()
    document.add_heading("Bias variance tradeoff", level=1)
    document.add_paragraph("Complex models can reduce bias but increase variance.")
    document.save(docx_path)

    segments, meta = extract_material(docx_path)

    assert meta["status"] == "indexed"
    assert segments[0]["locator_type"] == "heading"
    assert segments[0]["locator"] == "Bias variance tradeoff"


def test_xlsx_extraction_preserves_sheet_locator(tmp_path: Path) -> None:
    from openpyxl import Workbook

    xlsx_path = tmp_path / "dictionary.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Variables"
    sheet.append(["column", "meaning"])
    sheet.append(["sales", "weekly sales"])
    workbook.save(xlsx_path)

    segments, meta = extract_material(xlsx_path)

    assert meta["status"] == "indexed"
    assert segments[0]["locator_type"] == "sheet"
    assert segments[0]["locator"] == "Variables"


def test_tex_extraction_preserves_question_locators(tmp_path: Path) -> None:
    tex_path = tmp_path / "mock.tex"
    tex_path.write_text(
        r"""
\section*{Regression}
\begin{enumerate}
\item Explain training and test error.
\item Calculate the prediction error for a simple model.
\end{enumerate}
""",
        encoding="utf-8",
    )

    segments, meta = extract_material(tex_path)

    assert meta["status"] == "indexed"
    assert any(segment["locator_type"] == "question" for segment in segments)
    assert any("training and test error" in segment["text"] for segment in segments)


def test_csv_extraction_indexes_schema_and_sample(tmp_path: Path) -> None:
    csv_path = tmp_path / "customers.csv"
    csv_path.write_text("customer_id,spend\n1,10\n2,20\n", encoding="utf-8")

    segments, meta = extract_material(csv_path, config={"csv_max_rows": 2})

    assert meta["method"] == "csv_summary"
    assert segments[0]["locator_type"] == "rows"
    assert "Columns (2): customer_id | spend" in segments[0]["text"]


def test_image_caption_flow_creates_stub_then_indexes_caption(tmp_path: Path) -> None:
    image_path = tmp_path / "tree.png"
    image_path.write_bytes(b"not-a-real-image-but-good-enough-for-caption-flow")
    caption_root = tmp_path / "captions"

    segments, meta = extract_material(image_path, caption_root=caption_root)

    assert segments == []
    assert meta["status"] == "caption_needed"
    caption_path = Path(meta["caption_path"])
    assert caption_path.exists()

    caption_path.write_text("Decision tree diagram showing root split and leaf nodes.", encoding="utf-8")
    segments, meta = extract_material(image_path, caption_root=caption_root)

    assert meta["status"] == "indexed"
    assert segments[0]["locator_type"] == "caption"
    assert "Decision tree diagram" in segments[0]["text"]


def test_caption_path_is_stable_for_absolute_and_relative_paths(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    caption_root = tmp_path / "captions"

    assert caption_path_for_image(image_path, caption_root) == caption_path_for_image(image_path.resolve(), caption_root)


def test_zip_extraction_indexes_manifest_and_supported_members(tmp_path: Path) -> None:
    zip_path = tmp_path / "bundle.zip"
    with zipfile.ZipFile(zip_path, "w") as archive:
        archive.writestr("notes.txt", "Cross validation balances bias and variance.")
        archive.writestr("data.csv", "x,y\n1,2\n")

    segments, meta = extract_material(zip_path, caption_root=tmp_path / "captions")

    assert meta["method"] == "zip"
    assert any(segment["locator_type"] == "archive" for segment in segments)
    assert any("Cross validation" in segment["text"] for segment in segments)
